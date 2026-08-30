"""Best-effort text extraction from stored project files.

Every extractor is wrapped so a parser failure returns ``("error", "")`` rather
than raising: the upload must never crash and parser internals must never reach
the client. Archives, images and CAD formats are opaque -- never opened here.
All extracted content is untrusted inert text.
"""

from __future__ import annotations

from pathlib import Path

_PLAINTEXT = {".txt", ".md", ".csv", ".tsv", ".json"}

_XLSX_MAX_SHEETS = 10
_XLSX_MAX_ROWS = 1000
_XLSX_MAX_COLS = 50


def _pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(path: Path) -> str:
    import docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def _xlsx(path: Path) -> str:
    import openpyxl

    # data_only=True returns the cached value, never evaluates a formula.
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        out: list[str] = []
        for sheet_index, ws in enumerate(wb.worksheets):
            if sheet_index >= _XLSX_MAX_SHEETS:
                break
            out.append(f"# {ws.title}")
            for row_index, row in enumerate(ws.iter_rows(values_only=True)):
                if row_index >= _XLSX_MAX_ROWS:
                    break
                cells = ["" if v is None else str(v) for v in row[:_XLSX_MAX_COLS]]
                out.append("\t".join(cells))
        return "\n".join(out)
    finally:
        wb.close()


def _pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def extract_text(path: Path, extension: str, max_bytes: int) -> tuple[str, str]:
    """Return ``(status, text)`` where status is one of
    ``ok | truncated | empty | unsupported | error``."""
    ext = extension.lower()
    try:
        if ext == ".pdf":
            text = _pdf(path)
        elif ext in _PLAINTEXT:
            text = Path(path).read_bytes().decode("utf-8", "replace")
        elif ext == ".docx":
            text = _docx(path)
        elif ext == ".xlsx":
            text = _xlsx(path)
        elif ext == ".pptx":
            text = _pptx(path)
        else:
            return ("unsupported", "")
    except Exception:
        return ("error", "")

    if not text or not text.strip():
        return ("empty", "")

    encoded = text.encode("utf-8")
    if len(encoded) > max_bytes:
        return ("truncated", encoded[:max_bytes].decode("utf-8", "ignore"))
    return ("ok", text)
