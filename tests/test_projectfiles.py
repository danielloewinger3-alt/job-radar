"""Project-file API: upload streaming + quotas, signature validation, text
extraction, safe downloads and ownership.

No network, isolated DB, project-file store redirected into tmp_path by the
conftest isolated_db fixture.
"""

import io
import zipfile
from contextlib import contextmanager
from pathlib import Path

import pytest

from backend import config
from backend.models import Project
from backend.projectfiles.models import ProjectFile

PNG_MAGIC = b"\x89PNG\r\n\x1a\n" + b"\x00" * 32
PDF_BYTES = b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj\n<<>>\nendobj\ntrailer<<>>\n%%EOF\n"


def _project(session, title="Proj"):
    p = Project(title=title)
    session.add(p)
    session.commit()
    session.refresh(p)
    return p


def _upload(client, project_id, *, name="doc.pdf", data=PDF_BYTES,
            content_type="application/octet-stream", description=None,
            ai_context_enabled=None):
    form = {}
    if description is not None:
        form["description"] = description
    if ai_context_enabled is not None:
        form["ai_context_enabled"] = ai_context_enabled
    return client.post(
        f"/api/projects/{project_id}/files",
        files={"file": (name, io.BytesIO(data), content_type)},
        data=form,
    )


def _store_files():
    return [
        p for p in Path(config.PROJECTFILES_DIR).iterdir() if p.is_file()
    ]


def _temp_files():
    return [
        p for p in Path(config.PROJECTFILES_DIR).iterdir()
        if p.name.startswith(".upload-")
    ]


def _zip_bytes(members):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for n, c in members.items():
            zf.writestr(n, c)
    return buf.getvalue()


def _docx_bytes(text="hello docx body"):
    import docx

    d = docx.Document()
    d.add_paragraph(text)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _xlsx_bytes():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "cellvalue"
    ws["A2"] = "=A1&A1"  # formula -- must never be evaluated
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes(text="hello pptx text"):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Contract shape + happy path
# --------------------------------------------------------------------------- #
def test_list_is_wrapped(client, db_session):
    pid = _project(db_session).id
    body = client.get(f"/api/projects/{pid}/files").json()
    assert body["schema_version"] == 1
    assert body["files"] == []


def test_list_unknown_project_404(client):
    assert client.get("/api/projects/999/files").status_code == 404


def test_upload_happy_path(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="spec.pdf", data=PDF_BYTES)
    assert r.status_code == 201
    f = r.json()["file"]
    assert f["original_name"] == "spec.pdf"
    assert f["extension"] == ".pdf"
    assert f["byte_size"] == len(PDF_BYTES)
    assert len(f["sha256"]) == 64
    row = db_session.exec(
        __import__("sqlmodel").select(ProjectFile)
    ).one()
    assert row.stored_name.endswith(".pdf") and "/" not in row.stored_name
    assert row.stored_name != "spec.pdf"


# --------------------------------------------------------------------------- #
# Status codes
# --------------------------------------------------------------------------- #
def test_disallowed_extension_is_415(client, db_session):
    pid = _project(db_session).id
    assert _upload(client, pid, name="a.exe", data=b"MZ").status_code == 415


def test_signature_mismatch_is_415(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="fake.pdf", data=PNG_MAGIC)
    assert r.status_code == 415
    assert db_session.exec(__import__("sqlmodel").select(ProjectFile)).all() == []
    assert _store_files() == []


def test_individual_size_limit_is_413(client, db_session, monkeypatch):
    pid = _project(db_session).id
    monkeypatch.setattr(config, "MAX_PROJECT_FILE_BYTES", 8)
    r = _upload(client, pid, name="big.pdf", data=PDF_BYTES)
    assert r.status_code == 413
    assert _temp_files() == []


def test_per_project_byte_limit_is_413(client, db_session, monkeypatch):
    pid = _project(db_session).id
    _upload(client, pid, name="one.pdf", data=PDF_BYTES)
    monkeypatch.setattr(config, "MAX_PROJECT_FILES_PER_PROJECT_BYTES", len(PDF_BYTES) + 1)
    assert _upload(client, pid, name="two.pdf", data=PDF_BYTES).status_code == 413


def test_global_byte_limit_is_413(client, db_session, monkeypatch):
    pid = _project(db_session).id
    _upload(client, pid, name="one.pdf", data=PDF_BYTES)
    monkeypatch.setattr(config, "MAX_PROJECT_FILES_TOTAL_BYTES", len(PDF_BYTES) + 1)
    assert _upload(client, pid, name="two.pdf", data=PDF_BYTES).status_code == 413


def test_file_count_limit_is_409(client, db_session, monkeypatch):
    pid = _project(db_session).id
    monkeypatch.setattr(config, "MAX_PROJECT_FILES_PER_PROJECT", 1)
    assert _upload(client, pid, name="a.pdf").status_code == 201
    r = _upload(client, pid, name="b.pdf")
    assert r.status_code == 409
    assert _temp_files() == []


@pytest.mark.parametrize("bad", ["../../etc/passwd", "..\\..\\x.pdf", "..", "."])
def test_malformed_filename_is_400(client, db_session, bad):
    pid = _project(db_session).id
    assert _upload(client, pid, name=bad).status_code == 400


def test_safe_filename_with_internal_dots_accepted(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="my..notes.final.pdf", data=PDF_BYTES)
    assert r.status_code == 201
    assert r.json()["file"]["original_name"] == "my..notes.final.pdf"


def test_malformed_ai_context_enabled_is_400(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="a.pdf", ai_context_enabled="maybe")
    assert r.status_code == 400


# --------------------------------------------------------------------------- #
# Disguised Office file + cleanup on failure
# --------------------------------------------------------------------------- #
def test_disguised_docx_stored_opaque(client, db_session):
    pid = _project(db_session).id
    fake = _zip_bytes({"foo.txt": "not an ooxml"})
    r = _upload(client, pid, name="fake.docx", data=fake)
    assert r.status_code == 201
    f = r.json()["file"]
    assert f["extract_status"] == "error"
    assert f["ai_context_enabled"] is False
    assert "traceback" not in r.text.lower()


def test_streaming_overflow_leaves_no_temp(client, db_session, monkeypatch):
    pid = _project(db_session).id
    monkeypatch.setattr(config, "MAX_PROJECT_FILE_BYTES", 4)
    _upload(client, pid, name="big.pdf", data=PDF_BYTES)
    assert _store_files() == []


def test_db_commit_failure_cleans_blob(client, db_session, monkeypatch):
    pid = _project(db_session).id
    import backend.projectfiles.router as pr

    real_get_session = pr.get_session

    @contextmanager
    def flaky():
        with real_get_session() as s:
            def boom():
                raise RuntimeError("commit boom")

            s.commit = boom
            yield s

    monkeypatch.setattr(pr, "get_session", flaky)
    r = _upload(client, pid, name="x.pdf", data=PDF_BYTES)
    assert r.status_code == 500
    assert "boom" not in r.text  # controlled body, no exception text
    # (do NOT call monkeypatch.undo(): the shared fixture also holds isolated_db's
    # PROJECTFILES_DIR / engine redirects.)
    db_session.expire_all()
    assert db_session.exec(__import__("sqlmodel").select(ProjectFile)).all() == []
    assert _store_files() == []


def test_recheck_helper_serialises_count_limit(client, db_session, monkeypatch):
    """Unit-level proof of the recheck-before-persist guard."""
    import backend.projectfiles.router as pr
    from backend.db import get_session

    pid = _project(db_session).id
    _upload(client, pid, name="a.pdf")
    monkeypatch.setattr(config, "MAX_PROJECT_FILES_PER_PROJECT", 1)
    with get_session() as s:
        with pytest.raises(Exception) as ei:
            pr._recheck_limits(s, pid, 10)
    assert ei.value.status_code == 409


# --------------------------------------------------------------------------- #
# Extraction
# --------------------------------------------------------------------------- #
def test_docx_extraction(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="d.docx", data=_docx_bytes("secret docx phrase"))
    fid = r.json()["file"]["file_id"]
    assert r.json()["file"]["extract_status"] == "ok"
    got = client.get(f"/api/projects/{pid}/files/{fid}").json()
    assert "secret docx phrase" in got["extracted_text"]


def test_xlsx_extraction_does_not_evaluate_formulas(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="s.xlsx", data=_xlsx_bytes())
    fid = r.json()["file"]["file_id"]
    text = client.get(f"/api/projects/{pid}/files/{fid}").json()["extracted_text"]
    assert "cellvalue" in text
    assert "=A1" not in text  # formula neither shown nor evaluated


def test_pptx_extraction(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="p.pptx", data=_pptx_bytes("deck phrase here"))
    fid = r.json()["file"]["file_id"]
    text = client.get(f"/api/projects/{pid}/files/{fid}").json()["extracted_text"]
    assert "deck phrase here" in text


def test_text_extract_byte_cap_truncates(client, db_session, monkeypatch):
    pid = _project(db_session).id
    monkeypatch.setattr(config, "PROJECT_FILE_TEXT_EXTRACT_MAX_BYTES", 10)
    r = _upload(client, pid, name="big.txt", data=b"x" * 500)
    fid = r.json()["file"]["file_id"]
    assert r.json()["file"]["extract_status"] == "truncated"
    got = client.get(f"/api/projects/{pid}/files/{fid}").json()
    assert len(got["extracted_text"].encode("utf-8")) <= 10


def test_corrupt_docx_is_error_not_crash(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="broken.docx", data=b"PK\x03\x04" + b"\x00" * 64)
    assert r.status_code == 201
    f = r.json()["file"]
    assert f["extract_status"] == "error"
    assert f["ai_context_enabled"] is False


@pytest.mark.parametrize("name,data", [
    ("part.step", b"ISO-10303-21;\nHEADER;"),
    ("model.stl", b"solid cube\nendsolid cube\n"),
    ("bundle.zip", None),
])
def test_cad_and_archive_are_opaque(client, db_session, name, data):
    pid = _project(db_session).id
    if data is None:
        data = _zip_bytes({"a.txt": "x"})
    r = _upload(client, pid, name=name, data=data)
    assert r.status_code == 201
    fid = r.json()["file"]["file_id"]
    assert r.json()["file"]["extract_status"] == "unsupported"
    # blob stored unchanged
    row = db_session.get(ProjectFile, fid)
    blob = (Path(config.PROJECTFILES_DIR) / row.stored_name).read_bytes()
    assert blob == data
    # cannot enable AI context
    assert client.patch(
        f"/api/projects/{pid}/files/{fid}", json={"ai_context_enabled": True}
    ).status_code == 409


# --------------------------------------------------------------------------- #
# Downloads
# --------------------------------------------------------------------------- #
def test_download_headers(client, db_session):
    pid = _project(db_session).id
    fid = _upload(client, pid, name="doc.pdf", data=PDF_BYTES).json()["file"]["file_id"]
    r = client.get(f"/api/projects/{pid}/files/{fid}/download")
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/octet-stream"
    assert r.headers["x-content-type-options"] == "nosniff"
    assert r.headers["content-disposition"].startswith("attachment;")
    assert r.content == PDF_BYTES


def test_download_non_ascii_filename(client, db_session):
    pid = _project(db_session).id
    fid = _upload(client, pid, name="café.pdf", data=PDF_BYTES).json()["file"]["file_id"]
    cd = client.get(f"/api/projects/{pid}/files/{fid}/download").headers["content-disposition"]
    assert "filename=\"caf.pdf\"" in cd
    assert "filename*=UTF-8''caf%C3%A9.pdf" in cd
    assert "\r" not in cd and "\n" not in cd


def test_download_unknown_file_404(client, db_session):
    pid = _project(db_session).id
    assert client.get(f"/api/projects/{pid}/files/999/download").status_code == 404


# --------------------------------------------------------------------------- #
# Ownership
# --------------------------------------------------------------------------- #
@pytest.mark.parametrize("verb", ["meta", "text", "download", "patch", "delete"])
def test_cross_project_access_is_404(client, db_session, verb):
    a = _project(db_session, "A").id
    b = _project(db_session, "B").id
    fid = _upload(client, a, name="doc.pdf", data=PDF_BYTES).json()["file"]["file_id"]
    if verb in ("meta", "text"):
        assert client.get(f"/api/projects/{b}/files/{fid}").status_code == 404
    elif verb == "download":
        assert client.get(f"/api/projects/{b}/files/{fid}/download").status_code == 404
    elif verb == "patch":
        assert client.patch(
            f"/api/projects/{b}/files/{fid}", json={"description": "x"}
        ).status_code == 404
    else:
        assert client.delete(f"/api/projects/{b}/files/{fid}").status_code == 404


# --------------------------------------------------------------------------- #
# PATCH / DELETE
# --------------------------------------------------------------------------- #
def test_patch_description_and_ai_context(client, db_session):
    pid = _project(db_session).id
    fid = _upload(client, pid, name="d.docx", data=_docx_bytes("body")).json()["file"]["file_id"]
    r = client.patch(
        f"/api/projects/{pid}/files/{fid}",
        json={"description": "notes", "ai_context_enabled": True},
    )
    assert r.status_code == 200
    assert r.json()["file"]["description"] == "notes"
    assert r.json()["file"]["ai_context_enabled"] is True


def test_patch_ai_context_on_nonreadable_is_409(client, db_session):
    pid = _project(db_session).id
    fid = _upload(
        client, pid, name="pic.png", data=PNG_MAGIC, content_type="image/png"
    ).json()["file"]["file_id"]
    assert client.patch(
        f"/api/projects/{pid}/files/{fid}", json={"ai_context_enabled": True}
    ).status_code == 409


def test_upload_requesting_ai_on_nonreadable_forced_false(client, db_session):
    pid = _project(db_session).id
    r = _upload(client, pid, name="pic.png", data=PNG_MAGIC,
                content_type="image/png", ai_context_enabled="true")
    assert r.status_code == 201
    assert r.json()["file"]["ai_context_enabled"] is False


def test_delete_removes_row_and_blob(client, db_session):
    pid = _project(db_session).id
    up = _upload(client, pid, name="doc.pdf", data=PDF_BYTES).json()["file"]
    fid = up["file_id"]
    row = db_session.get(ProjectFile, fid)
    blob = Path(config.PROJECTFILES_DIR) / row.stored_name
    assert blob.is_file()
    assert client.delete(f"/api/projects/{pid}/files/{fid}").json()["deleted"] is True
    assert not blob.exists()
    db_session.expire_all()
    assert db_session.get(ProjectFile, fid) is None
